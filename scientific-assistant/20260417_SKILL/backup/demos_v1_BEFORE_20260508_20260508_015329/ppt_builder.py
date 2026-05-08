"""
demos_v1/ppt_builder.py - MD → .pptx 결정론적 변환기 (LLM 없이)

사용자가 MD 파일/텍스트를 입력하면 python-pptx 로 직접 렌더링.
LLM 이 끼어들지 않으므로 코드 에러·토큰 제한 등 문제 없음.

MD 규칙:
- # 제목       → 제목 슬라이드 (부제는 다음 줄의 > 텍스트)
- ## 슬라이드  → 내용 슬라이드 시작
- - 불릿       → 본문 불릿 (들여쓰기로 level 조절)
- | 표 |       → 표 슬라이드 (자동 감지)
- ```lang\ncode\n``` → 코드 슬라이드 (모노스페이스)
- ![캡션](경로) → 이미지 슬라이드 (Phase 2)
- --- 또는 ## 재등장 → 슬라이드 분리
"""
import io
import os
import re
import time
import uuid
from typing import Any

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Cm, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False


# ============================================================
# 한글 폰트 폴백 체인
# ============================================================
KOREAN_FONT_CHAIN = ["맑은 고딕", "Malgun Gothic", "NanumGothic",
                    "Noto Sans KR", "Pretendard", "Arial"]
CODE_FONT_CHAIN = ["Consolas", "D2Coding", "Courier New", "Monaco"]


# ============================================================
# 테마 프리셋 — 템플릿명 → (accent, bg, text, code_bg, code_text, header_bg)
# ============================================================
THEME_PRESETS = {
    "minimal":   {"accent": (0xEA, 0xB3, 0x08), "header_bg": (0xF5, 0xF5, 0xF5),
                  "header_text": (0x22, 0x22, 0x22), "zebra": (0xFA, 0xFA, 0xFA),
                  "code_bg": (0x1E, 0x1E, 0x1E), "code_text": (0xDC, 0xDC, 0xDC)},
    "corporate": {"accent": (0x1F, 0x6F, 0xEB), "header_bg": (0x1F, 0x6F, 0xEB),
                  "header_text": (0xFF, 0xFF, 0xFF), "zebra": (0xEE, 0xF4, 0xFC),
                  "code_bg": (0x1E, 0x1E, 0x1E), "code_text": (0xDC, 0xDC, 0xDC)},
    "academic":  {"accent": (0x8B, 0x1A, 0x1A), "header_bg": (0x8B, 0x1A, 0x1A),
                  "header_text": (0xFF, 0xFF, 0xFF), "zebra": (0xF7, 0xEE, 0xEE),
                  "code_bg": (0x2B, 0x2B, 0x2B), "code_text": (0xE0, 0xE0, 0xE0)},
    "creative":  {"accent": (0xE0, 0x3E, 0x96), "header_bg": (0xE0, 0x3E, 0x96),
                  "header_text": (0xFF, 0xFF, 0xFF), "zebra": (0xFD, 0xEC, 0xF5),
                  "code_bg": (0x29, 0x1A, 0x36), "code_text": (0xF5, 0xE6, 0xF0)},
    "dark":      {"accent": (0x22, 0xC5, 0x5E), "header_bg": (0x22, 0xC5, 0x5E),
                  "header_text": (0x0A, 0x0A, 0x0A), "zebra": (0x1F, 0x1F, 0x2E),
                  "code_bg": (0x0A, 0x0A, 0x0A), "code_text": (0x9C, 0xDC, 0xFE)},
}
DEFAULT_THEME = "minimal"


# ============================================================
# MD Parser
# ============================================================
def parse_md_to_outline(md_text: str) -> dict:
    """MD 텍스트 → outline dict 변환.

    Returns:
        {
            "meta": {"title": str, "subtitle": str},
            "slides": [
                {"type": "title", "title": str, "subtitle": str},
                {"type": "content", "title": str, "bullets": [{"text", "level"}]},
                {"type": "table", "title": str, "headers": [...], "rows": [[...]]},
                {"type": "code", "title": str, "language": str, "code": str},
                ...
            ]
        }
    """
    lines = md_text.replace("\r\n", "\n").split("\n")
    meta = {"title": "", "subtitle": ""}
    slides: list[dict] = []
    current: dict | None = None  # 현재 빌드 중인 슬라이드
    in_code_block = False
    code_lang = ""
    code_buffer: list[str] = []
    table_buffer: list[str] = []
    in_table = False

    def flush_table():
        """누적된 table_buffer 를 table 슬라이드로 변환 + current 로 승격."""
        nonlocal current, table_buffer, in_table
        if not table_buffer:
            return
        # header | separator | rows
        rows_raw = [l for l in table_buffer if l.strip().startswith("|")]
        if len(rows_raw) < 2:
            table_buffer = []
            in_table = False
            return

        def _split_row(row: str) -> list[str]:
            row = row.strip().strip("|")
            return [c.strip() for c in row.split("|")]

        headers = _split_row(rows_raw[0])
        # rows_raw[1] 은 구분자 (|---|---|) → skip
        body_rows = [_split_row(r) for r in rows_raw[2:] if r.strip()]
        title = current.get("title", "") if current else ""
        table_slide = {"type": "table", "title": title,
                       "headers": headers, "rows": body_rows}
        # 기존 current 가 content slide 였고 bullet 없으면 table 로 대체
        if current and current.get("type") == "content" and not current.get("bullets"):
            slides[-1] = table_slide
            current = table_slide
        else:
            slides.append(table_slide)
            current = table_slide
        table_buffer = []
        in_table = False

    for raw in lines:
        line = raw.rstrip()

        # 코드 블록 내부
        if in_code_block:
            if line.startswith("```"):
                # 코드 블록 종료 → 슬라이드로 확정
                code_slide = {
                    "type": "code",
                    "title": current.get("title", "") if current else "코드",
                    "language": code_lang,
                    "code": "\n".join(code_buffer),
                }
                if current and current.get("type") == "content" and not current.get("bullets"):
                    slides[-1] = code_slide
                else:
                    slides.append(code_slide)
                current = code_slide
                in_code_block = False
                code_lang = ""
                code_buffer = []
            else:
                code_buffer.append(raw)  # raw 쓰기 (들여쓰기 유지)
            continue

        # 테이블 버퍼링
        if line.strip().startswith("|"):
            in_table = True
            table_buffer.append(line)
            continue
        elif in_table:
            flush_table()

        # 코드 블록 시작
        if line.startswith("```"):
            in_code_block = True
            code_lang = line[3:].strip()
            code_buffer = []
            continue

        # --- 구분자 → 섹션 슬라이드 (다음 ## 가 제목으로 들어감)
        if line.strip() in ("---", "***", "==="):
            # 다음 ## 슬라이드 제목을 섹션으로 승격하기 위한 플래그
            slides.append({"type": "section", "title": "", "_pending_title": True})
            current = slides[-1]
            continue

        # # 제목 (문서 메타)
        m1 = re.match(r"^#\s+(.+)$", line)
        if m1:
            title_text = m1.group(1).strip()
            if not slides:
                meta["title"] = title_text
                slide = {"type": "title", "title": title_text, "subtitle": ""}
                slides.append(slide)
                current = slide
            else:
                slide = {"type": "title", "title": title_text, "subtitle": ""}
                slides.append(slide)
                current = slide
            continue

        # ## 슬라이드 시작
        m2 = re.match(r"^##\s+(.+)$", line)
        if m2:
            t = m2.group(1).strip()
            # 직전이 빈 섹션 슬라이드면 제목만 채우고 종료
            if current and current.get("type") == "section" \
                    and current.get("_pending_title"):
                current["title"] = t
                current.pop("_pending_title", None)
                current = None  # 섹션 직후 본문 자동 안 만듦
                continue
            slide = {"type": "content", "title": t, "bullets": []}
            slides.append(slide)
            current = slide
            continue

        # ### ~ ###### 는 하위 제목으로 쓰거나 불릿 level 로
        m3 = re.match(r"^(#{3,6})\s+(.+)$", line)
        if m3:
            # 현재 슬라이드에 sub-heading 불릿으로 추가
            if current and current.get("type") == "content":
                level = len(m3.group(1)) - 2  # ### → 1, #### → 2
                current["bullets"].append({"text": m3.group(2).strip(), "level": level, "bold": True})
            continue

        # > 인용 (title 슬라이드의 subtitle 또는 본문 인용)
        mq = re.match(r"^>\s?(.*)$", line)
        if mq and current:
            if current.get("type") == "title" and not current.get("subtitle"):
                current["subtitle"] = mq.group(1).strip()
            elif current.get("type") == "content":
                current["bullets"].append({"text": "“" + mq.group(1).strip() + "”",
                                           "level": 0, "italic": True})
            continue

        # - 불릿 / 1. 번호 매김
        mb = re.match(r"^(\s*)([-*+]|\d+\.)\s+(.+)$", line)
        if mb:
            indent = len(mb.group(1))
            level = indent // 2   # 2 space = 1 level
            text = mb.group(3).strip()
            if not current or current.get("type") != "content":
                # 상위 슬라이드 없으면 "내용" 슬라이드 자동 생성
                slide = {"type": "content", "title": "내용", "bullets": []}
                slides.append(slide)
                current = slide
            current.setdefault("bullets", []).append({"text": text, "level": level})
            continue

        # 빈 줄은 무시
        if not line.strip():
            continue

        # 그 외 일반 텍스트 → 현재 슬라이드 본문에 추가
        if current and current.get("type") == "content":
            current.setdefault("bullets", []).append({"text": line.strip(), "level": 0})
        elif current and current.get("type") == "title" and not current.get("subtitle"):
            current["subtitle"] = line.strip()

    # 루프 종료 — 남은 테이블 flush
    if in_table:
        flush_table()

    # 슬라이드 하나도 없으면 안내 슬라이드
    if not slides:
        slides.append({"type": "title", "title": "빈 문서",
                       "subtitle": "MD 내용이 비어 있습니다."})

    # 미완성 섹션 슬라이드(다음 ## 가 안 와서 제목 비어있음) 정리
    slides = [s for s in slides
              if not (s.get("type") == "section" and s.get("_pending_title"))]

    # 코드 슬라이드 제목 중복 → (part N/M) 자동 부여
    _dedupe_code_titles(slides)

    # Smart Layout — content 슬라이드를 내용 모양 보고 분기
    _apply_smart_layout(slides)

    return {"meta": meta, "slides": slides}


# ============================================================
# Smart Layout — content 슬라이드 자동 재분류
# ============================================================
_BIGNUM_RE = re.compile(
    r"^\*{0,2}([+-]?\d[\d,\.]*\s*[%℃$₩원개명배회시간분초년월일])\*{0,2}$"
)
_VS_RE = re.compile(r"\bvs\.?\b|대\s*비|⟷|↔", re.IGNORECASE)


def _strip_md(text: str) -> str:
    """MD 강조 기호 제거 (스마트 레이아웃 매칭용 — 표시할 때 쓰는 함수 아님)."""
    t = text.strip()
    t = re.sub(r"^[*_]{1,3}|[*_]{1,3}$", "", t)
    return t.strip()


def _is_bignumber(bullets: list[dict]) -> bool:
    """단독/짧은 bullet 이 숫자+단위 형태면 빅넘버 슬라이드로."""
    if not bullets or len(bullets) > 3:
        return False
    return bool(_BIGNUM_RE.match(_strip_md(bullets[0].get("text", ""))))


def _is_quote(bullets: list[dict]) -> bool:
    """첫 bullet 이 italic + 따옴표로 감싸진 인용 형태."""
    if not bullets:
        return False
    b = bullets[0]
    text = b.get("text", "").strip()
    return b.get("italic") and (text.startswith("“") or text.startswith('"'))


def _apply_smart_layout(slides: list[dict]) -> None:
    """content 슬라이드를 statement/compare2/grid/quote/bignumber 로 자동 변환."""
    for s in slides:
        if s.get("type") != "content":
            continue
        bullets = s.get("bullets", [])
        title = s.get("title", "")

        # 0) 빅넘버 — bullet 이 숫자+단위 형태
        if _is_bignumber(bullets):
            big = _strip_md(bullets[0]["text"])
            label = bullets[1]["text"] if len(bullets) > 1 else ""
            s.clear()
            s.update({"type": "bignumber", "title": title,
                      "number": big, "label": label})
            continue

        # 1) 인용 — 첫 bullet 이 italic + 따옴표
        if _is_quote(bullets):
            quote = bullets[0].get("text", "").strip().strip("“”\"")
            cite = bullets[1].get("text", "").strip() if len(bullets) > 1 else ""
            s.clear()
            s.update({"type": "quote", "title": title,
                      "quote": quote, "cite": cite})
            continue

        # 2) statement — bullet 0개 또는 1개 짧은 문장
        if len(bullets) == 0 or (len(bullets) == 1
                                 and len(bullets[0].get("text", "")) <= 60):
            stmt = bullets[0].get("text", "") if bullets else ""
            s.clear()
            s.update({"type": "statement", "title": title, "statement": stmt})
            continue

        # 3) compare2 — 정확히 2개 top-level bullet, 또는 제목에 "vs"
        top = [b for b in bullets if b.get("level", 0) == 0]
        if len(top) == 2 or (_VS_RE.search(title) and len(top) >= 2):
            left = top[0]
            right = top[1]
            # 각 항목의 자식 (level >= 1) 수집
            def children_of(idx):
                kids = []
                start = bullets.index(top[idx]) + 1
                for b in bullets[start:]:
                    if b.get("level", 0) == 0:
                        break
                    kids.append(b.get("text", ""))
                return kids
            s.clear()
            s.update({"type": "compare2", "title": title,
                      "left_title": left.get("text", ""),
                      "left_items": children_of(0),
                      "right_title": right.get("text", ""),
                      "right_items": children_of(1)})
            continue

        # 4) grid — 정확히 3~4개 top-level bullet
        if 3 <= len(top) <= 4:
            cards = []
            for i, t in enumerate(top):
                start = bullets.index(t) + 1
                desc_lines = []
                for b in bullets[start:]:
                    if b.get("level", 0) == 0:
                        break
                    desc_lines.append(b.get("text", ""))
                cards.append({"title": t.get("text", ""),
                              "desc": " · ".join(desc_lines)})
            s.clear()
            s.update({"type": "grid", "title": title, "cards": cards})
            continue

        # 5) 그 외 — 그대로 content (불릿)


def _dedupe_code_titles(slides: list[dict]) -> None:
    """같은 제목의 연속 코드 슬라이드에 (part N/M) 표시 — in-place."""
    n = len(slides)
    i = 0
    while i < n:
        s = slides[i]
        if s.get("type") != "code":
            i += 1
            continue
        # 같은 제목 + 연속 code 묶기
        title = s.get("title", "")
        j = i
        while j < n and slides[j].get("type") == "code" \
                and slides[j].get("title", "") == title:
            j += 1
        group_size = j - i
        if group_size > 1:
            for k, idx in enumerate(range(i, j), start=1):
                slides[idx]["title"] = f"{title} (part {k}/{group_size})"
        i = j


# ============================================================
# Renderer (python-pptx)
# ============================================================
def _apply_korean_font(run, bold: bool = False, italic: bool = False,
                      size_pt: int = 18, is_code: bool = False):
    """run 에 한글 친화 폰트 + 속성 적용."""
    font = run.font
    font.name = (CODE_FONT_CHAIN[0] if is_code else KOREAN_FONT_CHAIN[0])
    font.size = Pt(size_pt)
    font.bold = bold
    font.italic = italic


def _add_accent_bar(slide, theme: dict):
    """슬라이드 좌측에 얇은 컬러바 + 우측 상단 점 강조."""
    from pptx.enum.shapes import MSO_SHAPE
    accent = RGBColor(*theme["accent"])
    # 좌측 세로 바
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                  Inches(0.12), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    # 제목 아래 짧은 강조 라인
    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5),
                                       Inches(1.25), Inches(1.0), Inches(0.05))
    underline.fill.solid()
    underline.fill.fore_color.rgb = accent
    underline.line.fill.background()


def _draw_title_text(slide, text: str, theme: dict, size_pt: int = 28,
                      top_in: float = 0.45, left_in: float = 0.5,
                      width_in: float = 9.0, height_in: float = 0.75):
    """슬라이드에 명시적 제목 textbox 를 그림 (placeholder 없는 Blank 레이아웃 대응)."""
    if not text:
        return
    accent = RGBColor(*theme["accent"])
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in),
                                    Inches(width_in), Inches(height_in))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    _apply_korean_font(r, bold=True, size_pt=size_pt)
    r.font.color.rgb = accent
    return tb


def _add_title_slide(prs, slide_data: dict, theme: dict):
    from pptx.enum.shapes import MSO_SHAPE
    layout = _safe_layout(prs, 0)
    slide = prs.slides.add_slide(layout)
    accent = RGBColor(*theme["accent"])
    # 좌측 컬러 블록
    block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                   Inches(0.4), Inches(7.5))
    block.fill.solid(); block.fill.fore_color.rgb = accent
    block.line.fill.background()
    # 큰 타이틀 (가운데 정렬, 큰 글씨)
    title_text = slide_data.get("title", "")
    if title_text:
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.6),
                                       Inches(8.6), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title_text
        _apply_korean_font(r, bold=True, size_pt=44)
        r.font.color.rgb = accent
    # 서브타이틀
    subtitle_text = slide_data.get("subtitle", "")
    if subtitle_text:
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(4.8),
                                       Inches(8.6), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle_text
        _apply_korean_font(r, size_pt=22)
        r.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    return slide


def _add_section_slide(prs, slide_data: dict, theme: dict):
    """섹션 구분 슬라이드 — 큰 컬러 배경 + 제목만."""
    from pptx.enum.shapes import MSO_SHAPE
    layout = _safe_layout(prs, 6) or _safe_layout(prs, 5)  # Blank
    slide = prs.slides.add_slide(layout)
    accent = RGBColor(*theme["accent"])
    # 전체 배경 띠
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.5),
                                Inches(10), Inches(2.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = accent
    bg.line.fill.background()
    # 제목 텍스트
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.9),
                                  Inches(8.8), Inches(1.7))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = slide_data.get("title", "")
    _apply_korean_font(run, bold=True, size_pt=40)
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return slide


def _add_statement_slide(prs, slide_data: dict, theme: dict):
    """짧은 한 문장 → 가운데 큰 텍스트."""
    from pptx.enum.shapes import MSO_SHAPE
    layout = _safe_layout(prs, 6) or _safe_layout(prs, 5)
    slide = prs.slides.add_slide(layout)
    _add_accent_bar(slide, theme)
    accent = RGBColor(*theme["accent"])
    # 작은 제목 (위)
    title = slide_data.get("title", "")
    if title:
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.5),
                                       Inches(8.8), Inches(0.6))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = title
        _apply_korean_font(r, bold=True, size_pt=20)
        r.font.color.rgb = accent
    # 큰 본문 (가운데)
    stmt = slide_data.get("statement", "")
    if stmt:
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.6),
                                       Inches(8.4), Inches(2.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = stmt
        _apply_korean_font(r, bold=True, size_pt=40)
    return slide


def _add_quote_slide(prs, slide_data: dict, theme: dict):
    """인용 — 큰 따옴표 + 가운데 인용문 + 출처."""
    from pptx.enum.shapes import MSO_SHAPE
    layout = _safe_layout(prs, 6) or _safe_layout(prs, 5)
    slide = prs.slides.add_slide(layout)
    _add_accent_bar(slide, theme)
    accent = RGBColor(*theme["accent"])
    # 큰 따옴표 장식
    tb_q = slide.shapes.add_textbox(Inches(0.6), Inches(1.0),
                                     Inches(2.0), Inches(2.0))
    pq = tb_q.text_frame.paragraphs[0]
    rq = pq.add_run()
    rq.text = "\u201C"
    _apply_korean_font(rq, bold=True, size_pt=140)
    rq.font.color.rgb = accent
    # 인용 본문
    quote = slide_data.get("quote", "")
    tb = slide.shapes.add_textbox(Inches(1.5), Inches(2.4),
                                   Inches(7.0), Inches(3.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = quote
    _apply_korean_font(r, italic=True, size_pt=28)
    # 출처
    cite = slide_data.get("cite", "")
    if cite:
        tb_c = slide.shapes.add_textbox(Inches(1.5), Inches(5.8),
                                         Inches(7.0), Inches(0.6))
        pc = tb_c.text_frame.paragraphs[0]
        pc.alignment = PP_ALIGN.RIGHT
        rc = pc.add_run()
        rc.text = "— " + cite
        _apply_korean_font(rc, size_pt=16)
        rc.font.color.rgb = accent
    return slide


def _add_bignumber_slide(prs, slide_data: dict, theme: dict):
    """빅넘버 — 거대한 숫자 + 라벨."""
    layout = _safe_layout(prs, 6) or _safe_layout(prs, 5)
    slide = prs.slides.add_slide(layout)
    _add_accent_bar(slide, theme)
    accent = RGBColor(*theme["accent"])
    title = slide_data.get("title", "")
    if title:
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.5),
                                       Inches(8.8), Inches(0.6))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = title
        _apply_korean_font(r, bold=True, size_pt=20)
        r.font.color.rgb = accent
    # 거대한 숫자 (가운데)
    num = slide_data.get("number", "")
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.8),
                                   Inches(9.0), Inches(3.5))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = num
    _apply_korean_font(r, bold=True, size_pt=140)
    r.font.color.rgb = accent
    # 라벨
    label = slide_data.get("label", "")
    if label:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(5.6),
                                       Inches(9.0), Inches(1.0))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        _apply_korean_font(r, size_pt=24)
    return slide


def _add_compare2_slide(prs, slide_data: dict, theme: dict):
    """2단 비교 — 좌우 카드."""
    from pptx.enum.shapes import MSO_SHAPE
    layout = _safe_layout(prs, 5)
    slide = prs.slides.add_slide(layout)
    _add_accent_bar(slide, theme)
    accent = RGBColor(*theme["accent"])
    _draw_title_text(slide, slide_data.get("title", "") or "비교", theme,
                       size_pt=28)
    # 좌우 카드 2개
    card_w, card_h = Inches(4.3), Inches(4.8)
    gap = Inches(0.3)
    top = Inches(1.8)
    left_x = Inches(0.5)
    right_x = Inches(0.5) + card_w + gap

    def draw_card(x, head, items, color: RGBColor):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top,
                                       card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
        card.line.color.rgb = color
        card.line.width = Pt(2)
        # 헤드 띠
        head_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top,
                                           card_w, Inches(0.7))
        head_bar.fill.solid()
        head_bar.fill.fore_color.rgb = color
        head_bar.line.fill.background()
        tb_h = slide.shapes.add_textbox(x, top + Inches(0.1),
                                          card_w, Inches(0.6))
        ph = tb_h.text_frame.paragraphs[0]
        ph.alignment = PP_ALIGN.CENTER
        rh = ph.add_run()
        rh.text = head
        _apply_korean_font(rh, bold=True, size_pt=20)
        rh.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        # 본문
        tb = slide.shapes.add_textbox(x + Inches(0.2), top + Inches(0.9),
                                        card_w - Inches(0.4), card_h - Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, it in enumerate(items or [""]):
            if i == 0:
                p = tf.paragraphs[0]
                p.text = ""
            else:
                p = tf.add_paragraph()
            r = p.add_run()
            r.text = "• " + it if it else ""
            _apply_korean_font(r, size_pt=16)

    # 살짝 다른 톤 (오른쪽은 회색계열로)
    color_l = accent
    color_r = RGBColor(0x55, 0x55, 0x55)
    draw_card(left_x, slide_data.get("left_title", "A"),
              slide_data.get("left_items", []), color_l)
    draw_card(right_x, slide_data.get("right_title", "B"),
              slide_data.get("right_items", []), color_r)
    return slide


def _add_grid_slide(prs, slide_data: dict, theme: dict):
    """3~4개 카드 그리드."""
    from pptx.enum.shapes import MSO_SHAPE
    layout = _safe_layout(prs, 5)
    slide = prs.slides.add_slide(layout)
    _add_accent_bar(slide, theme)
    accent = RGBColor(*theme["accent"])
    _draw_title_text(slide, slide_data.get("title", ""), theme, size_pt=28)
    cards = slide_data.get("cards", [])
    n = len(cards)
    if n == 0:
        return slide
    # 3개면 1행, 4개면 2x2
    if n <= 3:
        cols, rows = n, 1
    else:
        cols, rows = 2, 2
    margin_x, margin_y, gap = Inches(0.5), Inches(1.8), Inches(0.25)
    avail_w = Inches(10) - margin_x * 2 - gap * (cols - 1)
    avail_h = Inches(5.5) - gap * (rows - 1)
    card_w = Emu(int(avail_w / cols))
    card_h = Emu(int(avail_h / rows))

    for idx, c in enumerate(cards):
        row = idx // cols
        col = idx % cols
        x = margin_x + (card_w + gap) * col
        y = margin_y + (card_h + gap) * row
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y,
                                       card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
        card.line.color.rgb = accent
        card.line.width = Pt(1.5)
        # 좌측 액센트 바
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y,
                                             Inches(0.1), card_h)
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent
        accent_bar.line.fill.background()
        # 카드 제목
        tb = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.2),
                                        card_w - Inches(0.4), Inches(0.7))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = c.get("title", "")
        _apply_korean_font(r, bold=True, size_pt=18)
        r.font.color.rgb = accent
        # 카드 설명
        desc = c.get("desc", "")
        if desc:
            tb = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.95),
                                            card_w - Inches(0.4),
                                            card_h - Inches(1.1))
            tf = tb.text_frame
            tf.word_wrap = True
            r = tf.paragraphs[0].add_run()
            r.text = desc
            _apply_korean_font(r, size_pt=13)
    return slide


def _hex_to_rgb(hex_str: str, default=(0x88, 0x88, 0x88)) -> RGBColor:
    """'#1F6FEB' 또는 '1F6FEB' → RGBColor."""
    if not hex_str:
        return RGBColor(*default)
    s = hex_str.lstrip("#").strip()
    if len(s) != 6:
        return RGBColor(*default)
    try:
        return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    except ValueError:
        return RGBColor(*default)


_SHAPE_MAP = {
    "rect":          "RECTANGLE",
    "rectangle":     "RECTANGLE",
    "rounded_rect":  "ROUNDED_RECTANGLE",
    "rounded":       "ROUNDED_RECTANGLE",
    "oval":          "OVAL",
    "circle":        "OVAL",
    "ellipse":       "OVAL",
    "triangle":      "ISOCELES_TRIANGLE",
    "diamond":       "DIAMOND",
    "pentagon":      "PENTAGON",
    "hexagon":       "HEXAGON",
    "star":          "STAR_5_POINT",
    "arrow_right":   "RIGHT_ARROW",
    "arrow_left":    "LEFT_ARROW",
    "arrow_up":      "UP_ARROW",
    "arrow_down":    "DOWN_ARROW",
    "chevron":       "CHEVRON",
    "callout":       "RECTANGULAR_CALLOUT",
    "cloud":         "CLOUD",
}


def _add_custom_slide(prs, slide_data: dict, theme: dict):
    """LLM이 생성한 도형 사양을 그대로 그리는 자유형 슬라이드.

    slide_data 스키마:
      {
        "type": "custom",
        "title": "선택",
        "background": "#FFFFFF" (선택, 슬라이드 전체 배경),
        "shapes": [
          {"shape": "circle", "x": 1.0, "y": 2.0, "w": 1.5, "h": 1.5,
           "fill": "#1F6FEB", "line": "#FFFFFF", "line_width": 1,
           "text": "1", "font_size": 24, "font_color": "#FFFFFF",
           "bold": true, "italic": false, "align": "center"},
          {"shape": "textbox", "x": 0.5, "y": 0.5, "w": 9, "h": 0.7,
           "text": "제목", "font_size": 28, "bold": true, "font_color": "#1F6FEB"},
          {"shape": "line", "x": 1, "y": 4, "w": 8, "h": 0,
           "line": "#888", "line_width": 2}
        ]
      }
    좌표 단위: inch (10x7.5 슬라이드 기준).
    """
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN as _ALIGN

    layout = _safe_layout(prs, 6) or _safe_layout(prs, 5) or _safe_layout(prs, 0)
    slide = prs.slides.add_slide(layout)
    accent = RGBColor(*theme["accent"])

    # 배경 (선택)
    bg_hex = slide_data.get("background")
    if bg_hex:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                     Inches(10), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = _hex_to_rgb(bg_hex, (0xFF, 0xFF, 0xFF))
        bg.line.fill.background()

    # 좌측 액센트 바 (옵션) — slide_data 에 "no_accent_bar": true 면 생략
    if not slide_data.get("no_accent_bar"):
        _add_accent_bar(slide, theme)

    # 상단 제목 (있을 때)
    title = slide_data.get("title", "")
    if title and not slide_data.get("no_title"):
        _draw_title_text(slide, title, theme, size_pt=24)

    align_map = {"left": _ALIGN.LEFT, "center": _ALIGN.CENTER,
                 "right": _ALIGN.RIGHT}

    for sp in slide_data.get("shapes", []):
        kind = (sp.get("shape") or "rect").lower()
        x = Inches(float(sp.get("x", 0)))
        y = Inches(float(sp.get("y", 0)))
        w = Inches(float(sp.get("w", 1)))
        h = Inches(float(sp.get("h", 1)))
        fill_hex = sp.get("fill")
        line_hex = sp.get("line")
        line_w = float(sp.get("line_width", 0))
        text = sp.get("text", "")
        font_size = int(sp.get("font_size", 14))
        font_color = sp.get("font_color")
        bold = bool(sp.get("bold", False))
        italic = bool(sp.get("italic", False))
        align = align_map.get(str(sp.get("align", "left")).lower(),
                               _ALIGN.LEFT)

        try:
            if kind == "textbox":
                # 순수 텍스트 박스 (도형 X) — 슬라이드 배경 흰색 가정 → 디폴트 검정
                tb = slide.shapes.add_textbox(x, y, w, h)
                tf = tb.text_frame
                tf.word_wrap = True
                if text:
                    p = tf.paragraphs[0]
                    p.alignment = align
                    r = p.add_run()
                    r.text = text
                    _apply_korean_font(r, bold=bold, italic=italic,
                                       size_pt=font_size)
                    # textbox 디폴트 = 검정 (흰 배경에서도 보임)
                    r.font.color.rgb = _hex_to_rgb(font_color, (0x22, 0x22, 0x22)) \
                                        if font_color else RGBColor(0x22, 0x22, 0x22)
            elif kind == "line":
                # 직선 (Connector)
                from pptx.util import Emu
                conn = slide.shapes.add_connector(1, x, y, x + w, y + h)
                conn.line.color.rgb = _hex_to_rgb(line_hex or "#888888")
                conn.line.width = Pt(max(1, line_w or 1.5))
            else:
                # 도형 (rect, circle, arrow, ...)
                ms_name = _SHAPE_MAP.get(kind, "RECTANGLE")
                ms = getattr(MSO_SHAPE, ms_name, MSO_SHAPE.RECTANGLE)
                shp = slide.shapes.add_shape(ms, x, y, w, h)
                # fill
                if fill_hex and fill_hex.lower() != "none":
                    shp.fill.solid()
                    shp.fill.fore_color.rgb = _hex_to_rgb(fill_hex,
                                                          theme["accent"])
                elif fill_hex == "none":
                    shp.fill.background()
                else:
                    shp.fill.solid()
                    shp.fill.fore_color.rgb = accent
                # line
                if line_hex:
                    shp.line.color.rgb = _hex_to_rgb(line_hex)
                    if line_w:
                        shp.line.width = Pt(line_w)
                else:
                    shp.line.fill.background()
                # text
                if text:
                    tf = shp.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.alignment = align
                    r = p.add_run()
                    r.text = text
                    _apply_korean_font(r, bold=bold, italic=italic,
                                       size_pt=font_size)
                    if font_color:
                        r.font.color.rgb = _hex_to_rgb(font_color,
                                                       (0x22, 0x22, 0x22))
                    else:
                        # 자동 명암 — 배경(fill)이 밝으면 검정, 어두우면 흰색
                        rgb_tuple = (0xFF, 0xFF, 0xFF)
                        if fill_hex and fill_hex.lower() != "none":
                            rgb_tuple = _hex_to_rgb(fill_hex,
                                                     (0xFF, 0xFF, 0xFF))
                            try:
                                rgb_tuple = (rgb_tuple[0], rgb_tuple[1],
                                             rgb_tuple[2])
                            except Exception:
                                pass
                        # 휘도 계산 (간이): R+G+B 평균
                        if not fill_hex or fill_hex.lower() == "none":
                            # 도형 fill 없음 → 액센트가 디폴트 (auto-fill)
                            avg = (theme["accent"][0] + theme["accent"][1]
                                   + theme["accent"][2]) // 3
                        else:
                            try:
                                fc = _hex_to_rgb(fill_hex, theme["accent"])
                                # RGBColor 는 hex 문자열로만 추출 가능 →
                                # _hex_to_rgb 결과를 그대로 사용
                                hexs = fill_hex.lstrip("#")
                                rr = int(hexs[0:2], 16)
                                gg = int(hexs[2:4], 16)
                                bb = int(hexs[4:6], 16)
                                avg = (rr + gg + bb) // 3
                            except Exception:
                                avg = 128
                        r.font.color.rgb = (RGBColor(0x22, 0x22, 0x22)
                                            if avg > 160
                                            else RGBColor(0xFF, 0xFF, 0xFF))
        except Exception as e:
            print(f"[ppt_builder] custom shape error ({kind}): {e}")
    return slide


def _add_content_slide(prs, slide_data: dict, theme: dict):
    layout = _safe_layout(prs, 1)
    slide = prs.slides.add_slide(layout)
    _add_accent_bar(slide, theme)
    _draw_title_text(slide, slide_data.get("title", ""), theme, size_pt=32)
    # 명시적 본문 textbox
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5),
                                    Inches(9.0), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    bullets = slide_data.get("bullets", [])
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = ""
        else:
            p = tf.add_paragraph()
        p.level = min(b.get("level", 0), 4)
        # 인덴트로 시각적 hierarchy
        prefix = "  " * p.level + "• "
        run = p.add_run()
        run.text = prefix + b.get("text", "")
        _apply_korean_font(run, bold=b.get("bold", False),
                          italic=b.get("italic", False), size_pt=18)
    return slide


def _add_table_slide(prs, slide_data: dict, theme: dict):
    layout = _safe_layout(prs, 5)
    slide = prs.slides.add_slide(layout)
    _add_accent_bar(slide, theme)
    _draw_title_text(slide, slide_data.get("title", "") or "표", theme,
                       size_pt=28)
    headers = slide_data.get("headers", [])
    rows = slide_data.get("rows", [])
    if not headers:
        return slide
    n_cols = len(headers)
    n_rows = len(rows) + 1
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(9.0)
    height = Inches(0.4 * n_rows + 0.2)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table
    header_bg = RGBColor(*theme["header_bg"])
    header_text = RGBColor(*theme["header_text"])
    zebra = RGBColor(*theme["zebra"])
    # 헤더
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = str(h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                _apply_korean_font(r, bold=True, size_pt=14)
                r.font.color.rgb = header_text
    # 본문 (zebra)
    for i, row in enumerate(rows, start=1):
        for j in range(n_cols):
            cell = tbl.cell(i, j)
            cell.text = str(row[j]) if j < len(row) else ""
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = zebra
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    _apply_korean_font(r, size_pt=12)
    return slide


def _add_code_slide(prs, slide_data: dict, theme: dict):
    from pptx.enum.shapes import MSO_SHAPE
    layout = _safe_layout(prs, 5)
    slide = prs.slides.add_slide(layout)
    _add_accent_bar(slide, theme)
    accent = RGBColor(*theme["accent"])
    lang = slide_data.get("language", "")
    title_text = (slide_data.get("title", "") or "코드") + (f" ({lang})" if lang else "")
    _draw_title_text(slide, title_text, theme, size_pt=24)
    code_bg = RGBColor(*theme["code_bg"])
    code_text = RGBColor(*theme["code_text"])
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9.0)
    height = Inches(5.5)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = code_bg
    box.line.color.rgb = accent
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.15)
    code = slide_data.get("code", "")
    lines = code.split("\n") or [""]
    for i, ln in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = ""
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = ln
        _apply_korean_font(run, size_pt=11, is_code=True)
        run.font.color.rgb = code_text
    return slide


def _purge_all_slides(prs):
    """템플릿에서 상속받은 기존 슬라이드를 파트·관계·목록 전부 완전히 제거.

    단순히 _sldIdLst 에서만 제거하면 실제 slide1.xml 파일은 패키지 안에 남아
    저장 시 'Duplicate name' 경고 + PowerPoint 오픈 실패가 발생.
    """
    slides = prs.slides
    sld_id_lst = slides._sldIdLst
    # 모든 슬라이드 id 요소 수집
    id_elems = list(sld_id_lst)
    pres_part = prs.part
    for sld_id in id_elems:
        rId = sld_id.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        # 1) sldIdLst 에서 제거
        sld_id_lst.remove(sld_id)
        if not rId:
            continue
        # 2) 관계 + 파트 완전 제거
        try:
            slide_part = pres_part.related_part(rId)
        except Exception:
            slide_part = None
        # 관계 drop
        try:
            pres_part.rels.pop(rId, None)
        except Exception:
            try:
                del pres_part.rels[rId]
            except Exception:
                pass
        # 패키지에서 파트 drop
        if slide_part is not None:
            try:
                slide_part.package.drop_part(slide_part)
            except Exception:
                # fallback: part 맵에서 제거
                try:
                    pkg = slide_part.package
                    parts_map = getattr(pkg, "_parts", None)
                    if parts_map is not None:
                        parts_map.pop(slide_part.partname, None)
                except Exception:
                    pass


def _safe_layout(prs, idx):
    """slide_layouts[idx] 를 안전하게 가져옴. 없으면 가능한 레이아웃 중 첫 번째."""
    try:
        return prs.slide_layouts[idx]
    except (IndexError, KeyError):
        layouts = list(prs.slide_layouts)
        return layouts[0] if layouts else None


def _make_presentation(template: str = DEFAULT_THEME):
    """python-pptx Presentation 생성.

    중요: 번들 템플릿(.pptx) 들이 'Blank' 레이아웃 하나뿐이고, 로드 후 기존
    슬라이드를 지우려 하면 파트 잔존 → 'Duplicate name' 경고 + PowerPoint 손상
    경고 발생. 디자인 정보(테마 색)는 어차피 우리가 코드로 다 그리므로,
    **빈 Presentation() 으로 시작** 이 가장 안전.
    """
    try:
        return Presentation()
    except Exception as e0:
        raise RuntimeError(
            "python-pptx 기본 Presentation() 생성 실패.\n"
            "해결: pip uninstall python-pptx -y && pip install python-pptx\n"
            f"원본 에러: {e0}"
        )


def render_outline_to_pptx(outline: dict, template: str = DEFAULT_THEME) -> bytes:
    """outline dict → .pptx bytes 반환."""
    if not _PPTX_AVAILABLE:
        raise RuntimeError("python-pptx 가 설치되지 않았습니다. pip install python-pptx")
    prs = _make_presentation(template)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    theme = THEME_PRESETS.get(template, THEME_PRESETS[DEFAULT_THEME])

    dispatch = {
        "title":     _add_title_slide,
        "section":   _add_section_slide,
        "statement": _add_statement_slide,
        "quote":     _add_quote_slide,
        "bignumber": _add_bignumber_slide,
        "compare2":  _add_compare2_slide,
        "grid":      _add_grid_slide,
        "content":   _add_content_slide,
        "table":     _add_table_slide,
        "code":      _add_code_slide,
        "custom":    _add_custom_slide,
    }
    for slide_data in outline.get("slides", []):
        fn = dispatch.get(slide_data.get("type", "content"), _add_content_slide)
        try:
            fn(prs, slide_data, theme)
        except Exception as e:
            err_slide = prs.slides.add_slide(_safe_layout(prs, 5))
            if err_slide.shapes.title:
                err_slide.shapes.title.text = f"[렌더링 실패] {slide_data.get('type', '?')}"
            print(f"[ppt_builder] slide render error: {e}")

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ============================================================
# 파일 저장 + ID 관리
# ============================================================
PPT_CACHE_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                             "..", "uploads", "ppt_cache")
os.makedirs(PPT_CACHE_DIR, exist_ok=True)


def save_pptx(pptx_bytes: bytes, hint: str = "presentation") -> dict:
    """.pptx bytes 를 임시 파일로 저장하고 id/경로 반환."""
    ppt_id = f"{int(time.time())}_{uuid.uuid4().hex[:8]}"
    safe_hint = re.sub(r"[^\w가-힣]+", "_", hint)[:40] or "presentation"
    fname = f"{safe_hint}_{ppt_id}.pptx"
    path = os.path.join(PPT_CACHE_DIR, fname)
    with open(path, "wb") as f:
        f.write(pptx_bytes)
    return {"id": ppt_id, "filename": fname, "path": path, "size": len(pptx_bytes)}


def get_pptx_path(ppt_id: str) -> str | None:
    """ID 로 .pptx 파일 경로 조회."""
    if not ppt_id or "/" in ppt_id or "\\" in ppt_id or ".." in ppt_id:
        return None
    for fname in os.listdir(PPT_CACHE_DIR):
        if ppt_id in fname and fname.endswith(".pptx"):
            return os.path.join(PPT_CACHE_DIR, fname)
    return None


# ============================================================
# 편의 함수: md_text 한 번에 .pptx 로
# ============================================================
def md_to_pptx_file(md_text: str, title_hint: str = "",
                    template: str = DEFAULT_THEME) -> dict:
    """MD 텍스트 → outline → .pptx 파일 저장 → id 반환."""
    outline = parse_md_to_outline(md_text)
    pptx_bytes = render_outline_to_pptx(outline, template=template)
    title = title_hint or outline.get("meta", {}).get("title", "presentation")
    info = save_pptx(pptx_bytes, hint=title)
    info["outline"] = outline
    info["slide_count"] = len(outline.get("slides", []))
    info["template"] = template
    return info


def list_templates() -> list[str]:
    """사용 가능한 테마명 목록."""
    return list(THEME_PRESETS.keys())
