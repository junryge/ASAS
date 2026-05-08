"""
demos_v1/ppt_templates/generate_templates.py
사용자 환경에서 1회 실행 → ppt_templates/*.pptx 5개 생성.

실행:
  cd scientific-assistant
  python demos_v1/ppt_templates/generate_templates.py

요구: python-pptx (pip install python-pptx)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# 5개 템플릿 색상 팔레트
THEMES = {
    "corporate": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "primary": RGBColor(0x1A, 0x73, 0xE8),      # 구글 블루
        "secondary": RGBColor(0x42, 0x85, 0xF4),
        "accent": RGBColor(0xFB, 0xBC, 0x04),
        "text": RGBColor(0x20, 0x20, 0x20),
        "muted": RGBColor(0x5F, 0x63, 0x68),
        "font_title": "맑은 고딕",
        "font_body": "맑은 고딕",
    },
    "academic": {
        "bg": RGBColor(0xF8, 0xF9, 0xFA),
        "primary": RGBColor(0x5F, 0x63, 0x68),
        "secondary": RGBColor(0x9A, 0xA0, 0xA6),
        "accent": RGBColor(0xD9, 0x3C, 0x25),
        "text": RGBColor(0x1F, 0x1F, 0x1F),
        "muted": RGBColor(0x70, 0x75, 0x7A),
        "font_title": "맑은 고딕",
        "font_body": "맑은 고딕",
    },
    "creative": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "primary": RGBColor(0xA8, 0x55, 0xF7),      # 보라
        "secondary": RGBColor(0xEC, 0x48, 0x99),    # 핑크
        "accent": RGBColor(0xFB, 0xBF, 0x24),
        "text": RGBColor(0x1F, 0x1F, 0x1F),
        "muted": RGBColor(0x6B, 0x72, 0x80),
        "font_title": "맑은 고딕",
        "font_body": "맑은 고딕",
    },
    "minimal": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "primary": RGBColor(0x00, 0x00, 0x00),
        "secondary": RGBColor(0x4B, 0x5C, 0x6B),
        "accent": RGBColor(0xEF, 0x44, 0x44),
        "text": RGBColor(0x0F, 0x0F, 0x0F),
        "muted": RGBColor(0x9C, 0xA3, 0xAF),
        "font_title": "맑은 고딕",
        "font_body": "맑은 고딕",
    },
    "dark": {
        "bg": RGBColor(0x0F, 0x14, 0x19),
        "primary": RGBColor(0x00, 0xD4, 0xFF),      # 네온 시안
        "secondary": RGBColor(0x00, 0xFF, 0x88),    # 네온 그린
        "accent": RGBColor(0xFF, 0x6B, 0x35),       # 오렌지
        "text": RGBColor(0xE5, 0xE7, 0xEB),
        "muted": RGBColor(0x9C, 0xA3, 0xAF),
        "font_title": "맑은 고딕",
        "font_body": "맑은 고딕",
    },
}


def _set_bg_color(slide, color):
    """슬라이드 배경색 설정."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_korean_font(run, font_name, size_pt, color, bold=False):
    """한글 폰트 설정 (east_asia 속성까지)."""
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    run.font.bold = bold
    # 한글 깨짐 방지: rPr 에 eastAsia 타이페이스 명시
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {'typeface': font_name})
        rPr.append(ea)
    else:
        ea.set('typeface', font_name)


def _add_text(slide, left, top, width, height, text, font, size, color, bold=False, align=PP_ALIGN.LEFT):
    """텍스트박스 추가."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_korean_font(run, font, size, color, bold)
    return box


def _add_rect(slide, left, top, width, height, color):
    """색상 사각형 (장식용)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def build_title_slide(prs, theme):
    """Layout 0: Title Slide (표지)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _set_bg_color(slide, theme["bg"])
    # 상단 액센트 바
    _add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.3), theme["primary"])
    # 제목 (중앙 큰글씨)
    _add_text(slide, Inches(1), Inches(2.5), Inches(11.33), Inches(1.5),
              "프레젠테이션 제목", theme["font_title"], 44, theme["text"], bold=True, align=PP_ALIGN.CENTER)
    # 부제
    _add_text(slide, Inches(1), Inches(4), Inches(11.33), Inches(0.6),
              "부제목 또는 요약", theme["font_body"], 20, theme["muted"], align=PP_ALIGN.CENTER)
    # 하단 바
    _add_rect(slide, Inches(0), prs.slide_height - Inches(0.3), prs.slide_width, Inches(0.3), theme["primary"])
    return slide


def build_section_header(prs, theme):
    """Layout 1: Section Header (섹션 구분)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 전체 배경을 primary 색으로
    _set_bg_color(slide, theme["primary"])
    # 큰 제목 중앙
    _add_text(slide, Inches(1), Inches(3), Inches(11.33), Inches(1.5),
              "섹션 제목", theme["font_title"], 54,
              RGBColor(0xFF, 0xFF, 0xFF) if sum(theme["primary"]) < 380 else RGBColor(0x20, 0x20, 0x20),
              bold=True, align=PP_ALIGN.CENTER)
    return slide


def build_content_slide(prs, theme):
    """Layout 2: Title and Content (일반 내용)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg_color(slide, theme["bg"])
    # 상단 제목
    _add_text(slide, Inches(0.7), Inches(0.4), Inches(12), Inches(0.8),
              "슬라이드 제목", theme["font_title"], 32, theme["primary"], bold=True)
    # 제목 밑 라인
    _add_rect(slide, Inches(0.7), Inches(1.3), Inches(2), Inches(0.06), theme["primary"])
    # 본문 영역 (placeholder 텍스트박스)
    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, sample in enumerate(["첫 번째 요점", "두 번째 요점", "세 번째 요점"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        run = p.add_run()
        run.text = f"• {sample}"
        _set_korean_font(run, theme["font_body"], 20, theme["text"])
    return slide


def build_two_column(prs, theme):
    """Layout 3: Two Content."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg_color(slide, theme["bg"])
    _add_text(slide, Inches(0.7), Inches(0.4), Inches(12), Inches(0.8),
              "두 단 비교", theme["font_title"], 32, theme["primary"], bold=True)
    _add_rect(slide, Inches(0.7), Inches(1.3), Inches(2), Inches(0.06), theme["primary"])
    # 가운데 구분선
    _add_rect(slide, Inches(6.65), Inches(1.7), Inches(0.03), Inches(5), theme["muted"])
    # 좌측
    _add_text(slide, Inches(0.7), Inches(1.6), Inches(5.8), Inches(0.5),
              "왼쪽 제목", theme["font_title"], 22, theme["secondary"], bold=True)
    _add_text(slide, Inches(0.7), Inches(2.2), Inches(5.8), Inches(4.5),
              "• 왼쪽 항목 1\n• 왼쪽 항목 2", theme["font_body"], 18, theme["text"])
    # 우측
    _add_text(slide, Inches(6.9), Inches(1.6), Inches(5.8), Inches(0.5),
              "오른쪽 제목", theme["font_title"], 22, theme["secondary"], bold=True)
    _add_text(slide, Inches(6.9), Inches(2.2), Inches(5.8), Inches(4.5),
              "• 오른쪽 항목 1\n• 오른쪽 항목 2", theme["font_body"], 18, theme["text"])
    return slide


def build_table_slide(prs, theme):
    """Layout 4: Title and Table (표)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg_color(slide, theme["bg"])
    _add_text(slide, Inches(0.7), Inches(0.4), Inches(12), Inches(0.8),
              "표 슬라이드", theme["font_title"], 32, theme["primary"], bold=True)
    _add_rect(slide, Inches(0.7), Inches(1.3), Inches(2), Inches(0.06), theme["primary"])
    # 샘플 표 (3x3)
    rows, cols = 3, 3
    left = Inches(0.7)
    top = Inches(2)
    width = Inches(12)
    height = Inches(3)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for c in range(cols):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = theme["primary"]
        cell.text_frame.text = f"컬럼 {c+1}"
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                _set_korean_font(run, theme["font_body"], 14, RGBColor(0xFF, 0xFF, 0xFF), bold=True)
    for r in range(1, rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text_frame.text = f"셀 {r},{c+1}"
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    _set_korean_font(run, theme["font_body"], 12, theme["text"])
    return slide


def build_title_only(prs, theme):
    """Layout 5: Title Only (이미지/차트 자유 배치용)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg_color(slide, theme["bg"])
    _add_text(slide, Inches(0.7), Inches(0.4), Inches(12), Inches(0.8),
              "이미지/차트 슬라이드", theme["font_title"], 32, theme["primary"], bold=True)
    _add_rect(slide, Inches(0.7), Inches(1.3), Inches(2), Inches(0.06), theme["primary"])
    # 자리표시용 placeholder 박스
    ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(2), Inches(9.33), Inches(5))
    ph.fill.solid()
    ph.fill.fore_color.rgb = theme["bg"]
    ph.line.color.rgb = theme["muted"]
    ph.line.width = Pt(1.5)
    tf = ph.text_frame
    tf.text = "이미지 또는 차트 삽입 영역"
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            _set_korean_font(run, theme["font_body"], 18, theme["muted"])
    return slide


def build_template(theme_name, theme, output_path):
    """주어진 테마로 6개 샘플 슬라이드 포함 .pptx 생성."""
    prs = Presentation()
    # 16:9 와이드스크린 (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_title_slide(prs, theme)
    build_section_header(prs, theme)
    build_content_slide(prs, theme)
    build_two_column(prs, theme)
    build_table_slide(prs, theme)
    build_title_only(prs, theme)

    prs.save(output_path)
    print(f"✅ {output_path} 생성 완료 ({theme_name})")


def main():
    out_dir = os.path.dirname(os.path.abspath(__file__))
    for name, theme in THEMES.items():
        build_template(name, theme, os.path.join(out_dir, f"{name}.pptx"))
    print()
    print(f"✨ 5개 템플릿 모두 생성됨: {out_dir}")
    print("   → PowerPoint 에서 열어서 디자인 확인/수정 가능")
    print("   → ppt_builder.py 가 이 파일들을 읽어 슬라이드 복제해 사용")


if __name__ == "__main__":
    main()
