"""주간보고 PPT 생성 스크립트 - 수정 금지"""
import sys
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


def _set_cell(table, row, col, text, bold=False, bg=None, align=PP_ALIGN.LEFT, size=10):
    cell = table.cell(row, col)
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = "맑은 고딕"
    p.alignment = align
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


def _add_text(slide, left, top, width, height, text, size=10, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = "맑은 고딕"
    p.alignment = align


def generate(data_path, output_path):
    with open(data_path, "r", encoding="utf-8") as f:
        projects = json.load(f)

    GRAY = RGBColor(0xD9, 0xD9, 0xD9)
    LGRAY = RGBColor(0xF2, 0xF2, 0xF2)
    RED = RGBColor(0xFF, 0x00, 0x00)
    YELLOW = RGBColor(0xFF, 0xD7, 0x00)
    BLUE = RGBColor(0x44, 0x72, 0xC4)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for idx, proj in enumerate(projects):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        _add_text(slide, Cm(1), Cm(0.5), Cm(20), Cm(1.2),
                  f"{idx+1}. {proj['name']}", size=24, bold=True)

        for i, c in enumerate([RED, YELLOW, BLUE]):
            bar = slide.shapes.add_shape(1, Cm(1 + 10.5 * i), Cm(2.0), Cm(10.5), Cm(0.25))
            bar.fill.solid()
            bar.fill.fore_color.rgb = c
            bar.line.fill.background()

        n_cur = len(proj.get("current", []))
        n_nxt = len(proj.get("next", []))
        data_rows = max(n_cur, n_nxt, 1)
        total_rows = 2 + data_rows + 1

        ts = slide.shapes.add_table(total_rows, 6, Cm(1), Cm(2.5), Cm(31.5), Cm(1.2 * total_rows))
        tbl = ts.table

        tbl.columns[0].width = Cm(12)
        tbl.columns[1].width = Cm(2.5)
        tbl.columns[2].width = Cm(2.5)
        tbl.columns[3].width = Cm(12)
        tbl.columns[4].width = Cm(2.5)
        tbl.columns[5].width = Cm(2.5)

        tbl.cell(0, 0).merge(tbl.cell(0, 2))
        tbl.cell(0, 3).merge(tbl.cell(0, 5))
        _set_cell(tbl, 0, 0, "금주 실적", bold=True, bg=GRAY, align=PP_ALIGN.CENTER, size=12)
        _set_cell(tbl, 0, 3, "차주 계획", bold=True, bg=GRAY, align=PP_ALIGN.CENTER, size=12)

        for ci, h in enumerate(["추진 내용", "납기", "진척율", "추진 내용", "납기", "진척율"]):
            _set_cell(tbl, 1, ci, h, bold=True, bg=LGRAY, align=PP_ALIGN.CENTER)

        for ri in range(data_rows):
            r = ri + 2
            if ri < n_cur:
                it = proj["current"][ri]
                _set_cell(tbl, r, 0, it.get("content", ""), size=9)
                _set_cell(tbl, r, 1, it.get("date", ""), align=PP_ALIGN.CENTER, size=9)
                _set_cell(tbl, r, 2, it.get("progress", ""), align=PP_ALIGN.CENTER, size=9)
            if ri < n_nxt:
                it = proj["next"][ri]
                _set_cell(tbl, r, 3, it.get("content", ""), size=9)
                _set_cell(tbl, r, 4, it.get("date", ""), align=PP_ALIGN.CENTER, size=9)
                _set_cell(tbl, r, 5, it.get("progress", ""), align=PP_ALIGN.CENTER, size=9)

        ir = 2 + data_rows
        tbl.cell(ir, 0).merge(tbl.cell(ir, 5))
        itxt = "Issue 및 협의사항"
        if proj.get("issues"):
            itxt = f"Issue 및 협의사항: {proj['issues']}"
        _set_cell(tbl, ir, 0, itxt, bold=True, bg=GRAY)

        _add_text(slide, Cm(1), Cm(17), Cm(20), Cm(0.8),
                  "● : 완료  ○ : 계획  ▶ : 진행중  ※ : Issue/특이사항", size=9)
        _add_text(slide, Cm(15), Cm(17.5), Cm(3), Cm(0.6),
                  str(idx + 1), size=10, align=PP_ALIGN.CENTER)

    prs.save(output_path)
    print(f"생성 완료: {output_path}")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("사용법: python gen_pptx.py <data.json> <output.pptx>")
        sys.exit(1)
    generate(sys.argv[1], sys.argv[2])
